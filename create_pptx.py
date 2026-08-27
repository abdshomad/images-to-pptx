# /// script
# requires-python = ">=3.9"
# dependencies = [
#     "python-pptx",
#     "pillow",
# ]
# ///

import sys
import os
import re
from pathlib import Path

# Configure UTF-8 stdout if supported
if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION ---
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif', '.tiff'}

# Colors (Modern Clean Aesthetic)
PRIMARY_COLOR = RGBColor(30, 41, 59)      # Slate 800
ACCENT_COLOR = RGBColor(37, 99, 235)      # Royal Blue 600
MUTED_TEXT = RGBColor(100, 116, 139)      # Slate 500

def natural_sort_key(s):
    """Sort strings with embedded numbers naturally (e.g., 1.jpg, 2.jpg, 10.jpg)."""
    return [int(text) if text.isdigit() else text.lower() for text in re.split(r'(\d+)', str(s))]

def clean_title(filename_stem: str) -> str:
    """Format filename into a clean, human-readable title."""
    cleaned = filename_stem.replace('_', ' ').replace('-', ' ')
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    if cleaned.isdigit():
        return f"Slide {cleaned}"
    return cleaned.title()

def create_presentation_for_folder(folder_path: Path, output_dir: Path):
    image_files = [
        f for f in folder_path.iterdir()
        if f.is_file() and f.suffix.lower() in IMAGE_EXTENSIONS
    ]
    
    if not image_files:
        print(f"Skipping '{folder_path.name}' (no images found)")
        return None
    
    image_files.sort(key=lambda f: natural_sort_key(f.name))
    print(f"Processing '{folder_path.name}' ({len(image_files)} images)...")
    
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # 1. COVER SLIDE
    # ----------------------------------------------------
    cover_slide = prs.slides.add_slide(blank_layout)
    
    # Accent top bar
    accent_bar = cover_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_WIDTH_INCHES), Inches(0.12)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_COLOR
    accent_bar.line.fill.background()
    
    # Title Box in Cover Slide
    title_box = cover_slide.shapes.add_textbox(
        Inches(1.5), Inches(2.4), Inches(10.333), Inches(2.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    folder_display_name = folder_path.name.replace('_', ' ').replace('-', ' ').title()
    p_title.text = folder_display_name
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.name = "Segoe UI"
    p_title.font.color.rgb = PRIMARY_COLOR
    
    p_sub = tf.add_paragraph()
    p_sub.text = f"{len(image_files)} Screenshots / Slides Presentation"
    p_sub.font.size = Pt(20)
    p_sub.font.name = "Segoe UI"
    p_sub.font.color.rgb = ACCENT_COLOR
    p_sub.space_before = Pt(14)
    
    # ----------------------------------------------------
    # 2. IMAGE SLIDES
    # ----------------------------------------------------
    margin_x = 0.8
    header_top = 0.5
    header_height = 0.75
    
    img_area_top = header_top + header_height + 0.15
    img_area_left = margin_x
    max_img_w = SLIDE_WIDTH_INCHES - (2 * margin_x)
    max_img_h = SLIDE_HEIGHT_INCHES - img_area_top - 0.45
    
    for idx, img_path in enumerate(image_files, start=1):
        slide = prs.slides.add_slide(blank_layout)
        
        # Header accent bar
        dot = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(margin_x), Inches(header_top + 0.05), Inches(0.08), Inches(0.35)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_COLOR
        dot.line.fill.background()
        
        # Slide Title text box
        title_box = slide.shapes.add_textbox(
            Inches(margin_x + 0.2), Inches(header_top), Inches(max_img_w - 0.2), Inches(header_height)
        )
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = Inches(0)
        tf_title.margin_top = Inches(0)
        
        p = tf_title.paragraphs[0]
        p.text = clean_title(img_path.stem)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = PRIMARY_COLOR
        
        p_index = tf_title.add_paragraph()
        p_index.text = f"Slide {idx} of {len(image_files)} - {img_path.name}"
        p_index.font.size = Pt(11)
        p_index.font.name = "Segoe UI"
        p_index.font.color.rgb = MUTED_TEXT
        p_index.space_before = Pt(2)
        
        # Image scaling & placement
        try:
            with Image.open(img_path) as pil_img:
                orig_w, orig_h = pil_img.size
                
            scale = min(max_img_w / orig_w, max_img_h / orig_h)
            placed_w = orig_w * scale
            placed_h = orig_h * scale
            
            placed_left = img_area_left + (max_img_w - placed_w) / 2
            placed_top = img_area_top + (max_img_h - placed_h) / 2
            
            slide.shapes.add_picture(
                str(img_path),
                Inches(placed_left),
                Inches(placed_top),
                width=Inches(placed_w),
                height=Inches(placed_h)
            )
        except Exception as e:
            print(f"Error loading '{img_path.name}': {e}")
            
    # Save PPTX
    out_file = output_dir / f"{folder_path.name}.pptx"
    prs.save(str(out_file))
    print(f"[SUCCESS] Generated: {out_file.name}")
    return out_file

def main():
    root_dir = Path(__file__).resolve().parent
    print(f"Scanning root directory: {root_dir}")
    
    subfolders = [p for p in root_dir.iterdir() if p.is_dir() and not p.name.startswith('.')]
    if not subfolders:
        print("No subdirectories found in current directory.")
        return
        
    created_count = 0
    for folder in sorted(subfolders, key=lambda f: natural_sort_key(f.name)):
        out = create_presentation_for_folder(folder, root_dir)
        if out:
            created_count += 1
            
    print(f"\nAll presentations created successfully! Total: {created_count} file(s).")

if __name__ == "__main__":
    main()
